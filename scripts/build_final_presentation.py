#!/usr/bin/env python3
"""Build Final Project — clean professional university-style deck.

Uses 'Simple Business Proposal' template for its clean white master
background and professional theme. Styled after the Team-115 reference:
blue header bars, white body, clean typography, no distracting artwork.
"""
import subprocess
import sys
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BLUE    = RGBColor(0x1B, 0x3A, 0x5C)
LBLUE   = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT  = RGBColor(0xE8, 0x6C, 0x00)
TEAL    = RGBColor(0x00, 0x7C, 0x6E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x2D, 0x2D, 0x2D)
BODY    = RGBColor(0x44, 0x44, 0x44)
LGRAY   = RGBColor(0x88, 0x88, 0x88)
LTBG    = RGBColor(0xF4, 0xF6, 0xF8)
TBL_H   = RGBColor(0x1B, 0x3A, 0x5C)
TBL_E   = RGBColor(0xEB, 0xEF, 0xF3)
TBL_O   = RGBColor(0xFB, 0xFB, 0xFB)
F = "Calibri"
nsP = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _del(prs):
    while len(prs.slides):
        rid = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rid)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def _tr(s, k="fade"):
    t = etree.SubElement(s._element, f"{{{nsP}}}transition")
    t.set("spd", "med"); t.set("advClick", "1")
    m = {"fade": "fade", "push": ("push", "l"), "wipe": ("wipe", "d"), "cover": ("cover", "l")}
    v = m.get(k, "fade")
    if isinstance(v, tuple):
        e = etree.SubElement(t, f"{{{nsP}}}{v[0]}"); e.set("dir", v[1])
    else:
        etree.SubElement(t, f"{{{nsP}}}{v}")


def _r(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    sh.line.fill.background(); return sh


def _t(s, l, t, w, h, txt, *, sz=14, c=BODY, b=False, a=PP_ALIGN.LEFT, va=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = va
    p = tf.paragraphs[0]; p.text = txt; p.font.name = F
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.alignment = a
    return tb


def _tm(s, l, t, w, h, lines, *, sz=13, c=BODY, sp=1.35, b=False, a=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.name = F; p.font.size = Pt(sz)
        p.font.color.rgb = c; p.font.bold = b; p.alignment = a
        p.line_spacing = sp; p.space_after = Pt(4)
    return tb


def _hdr(s, title, sw, *, sub=None):
    _r(s, 0, 0, sw, Inches(0.95), BLUE)
    _r(s, 0, Inches(0.95), sw, Inches(0.04), ACCENT)
    _t(s, Inches(0.6), Inches(0.15), sw - Inches(1.2), Inches(0.65),
       title, sz=24, c=WHITE, b=True, va=MSO_ANCHOR.MIDDLE)
    if sub:
        _t(s, sw - Inches(4.0), Inches(0.15), Inches(3.5), Inches(0.65),
           sub, sz=10, c=RGBColor(0xAA, 0xBB, 0xCC), a=PP_ALIGN.RIGHT, va=MSO_ANCHOR.MIDDLE)


def _footer(s, sw, sh, num):
    _r(s, 0, sh - Inches(0.3), sw, Inches(0.3), BLUE)
    _t(s, Inches(0.5), sh - Inches(0.28), Inches(4), Inches(0.24),
       "ICD-10 Code Prediction  ·  NLP Final Project",
       sz=8, c=RGBColor(0xAA, 0xBB, 0xCC), va=MSO_ANCHOR.MIDDLE)
    _t(s, sw - Inches(1.0), sh - Inches(0.28), Inches(0.6), Inches(0.24),
       str(num), sz=8, c=RGBColor(0xAA, 0xBB, 0xCC), a=PP_ALIGN.RIGHT, va=MSO_ANCHOR.MIDDLE)


def _tbl(s, rows, cw, left, top, *, best_col=None):
    nr, nc = len(rows), len(rows[0])
    g = s.shapes.add_table(nr, nc, left, top, sum(cw), Inches(0.36) * nr)
    tbl = g.table
    for ci, w in enumerate(cw): tbl.columns[ci].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = str(val)
            p = cell.text_frame.paragraphs[0]; p.font.name = F; p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_H
                p.font.color.rgb = WHITE; p.font.bold = True; p.font.size = Pt(11)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_E if ri % 2 == 0 else TBL_O
                p.font.color.rgb = DARK; p.font.size = Pt(11)
                if best_col and ci == best_col:
                    p.font.bold = True; p.font.color.rgb = TEAL
    return g


def _circ(s, x, y, n, *, bg=LBLUE, d=Inches(0.35)):
    dd = Inches(0.42) if int(n) >= 10 else d
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, dd, dd)
    c.fill.solid(); c.fill.fore_color.rgb = bg; c.line.fill.background()
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.text = str(n); p.font.name = F
    p.font.size = Pt(10) if int(n) >= 10 else Pt(12)
    p.font.bold = True; p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def _sec(s, label, x, y, *, c=LBLUE):
    _r(s, x, y, Inches(0.05), Inches(0.22), c)
    _t(s, x + Inches(0.12), y - Inches(0.02), Inches(4), Inches(0.25),
       label, sz=11, c=c, b=True)


def _img(s, path, left, top, *, width=None, height=None):
    p = Path(path)
    if p.exists():
        s.shapes.add_picture(str(p), left, top, width=width, height=height)
        return True
    return False


# ═══════════════════════════════════════════════════════════

def build(tmpl: Path, out: Path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    SW, SH = prs.slide_width, prs.slide_height
    BL = prs.slide_layouts[6]  # blank layout
    CL = Inches(0.6)
    CW = Inches(8.8)
    Y0 = Inches(1.2)
    ROOT = tmpl.parent

    # Refresh Notebook 01 EDA figures from embedded cell outputs (no screenshot crop)
    _ext = ROOT / "scripts" / "extract_notebook01_eda_plots.py"
    if _ext.exists():
        subprocess.run(
            [sys.executable, str(_ext)],
            cwd=str(ROOT),
            check=False,
        )

    # ═══ 1  TITLE ═══
    s = prs.slides.add_slide(BL)
    _r(s, 0, 0, SW, SH, BLUE)
    _r(s, 0, Inches(3.5), SW, Inches(0.04), ACCENT)
    _t(s, Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.6),
       "ICD-10 Code Prediction from\nMIMIC-IV Discharge Summaries",
       sz=32, c=WHITE, b=True, a=PP_ALIGN.LEFT)
    _t(s, Inches(0.8), Inches(2.8), Inches(8.4), Inches(0.35),
       "NLP Final Project  ·  Spring 2025", sz=14, c=RGBColor(0x99, 0xAA, 0xBB))
    _t(s, Inches(0.8), Inches(3.8), Inches(8.4), Inches(0.25),
       "Presented by", sz=11, c=RGBColor(0x88, 0x99, 0xAA))
    _t(s, Inches(0.8), Inches(4.15), Inches(8.4), Inches(0.35),
       "Devasai Sunder Tangella    ·    [Partner Name]",
       sz=16, c=WHITE, b=True)
    _t(s, Inches(0.8), Inches(4.65), Inches(5), Inches(0.22),
       "(Edit names in PowerPoint)", sz=9, c=RGBColor(0x66, 0x77, 0x88))
    _tr(s, "fade")

    # ═══ 2  AGENDA ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Agenda", SW)
    _footer(s, SW, SH, 2)
    _tr(s, "push")

    agenda = ["Problem Statement & Motivation", "Prior Work",
              "Dataset & Base Model (TF-IDF)", "Individual Contributions",
              "Proposed Approaches — Overview", "Approach 1: Chunk-BERT + Label Attention",
              "Approach 2: BiLSTM-LAAT", "Experimental Results",
              "Discussion & Challenges", "Future Work"]
    for ci, items in enumerate([agenda[:5], agenda[5:]]):
        xb = Inches(0.8) + ci * Inches(4.5)
        for i, item in enumerate(items):
            y = Y0 + Inches(0.15) + i * Inches(0.72)
            _circ(s, xb, y, i + 1 + ci * 5)
            _t(s, xb + Inches(0.55), y + Inches(0.02), Inches(3.7), Inches(0.32),
               item, sz=13, c=DARK, b=True)

    # ═══ 3  PROBLEM ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Problem Statement & Motivation", SW)
    _footer(s, SW, SH, 3)
    _tr(s, "wipe")

    _sec(s, "PROBLEM", CL, Y0)
    _tm(s, CL + Inches(0.15), Y0 + Inches(0.28), CW, Inches(0.8), [
        "Given a clinical discharge summary (~1,500 words), predict the set of",
        "ICD-10 diagnosis codes assigned to that hospital encounter.",
        "This is a multi-label classification task with 7,940 possible codes.",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "WHY IS THIS HARD?", CL, Y0 + Inches(1.1), c=ACCENT)
    _tm(s, CL + Inches(0.15), Y0 + Inches(1.38), CW, Inches(1.2), [
        "•   Massive label space: 7,940 unique ICD-10 codes after frequency filtering",
        "•   Extreme class imbalance: rare codes appear in <0.1% of admissions",
        "•   Long documents: avg ~1,500 words exceeds BERT's 512-token limit",
        "•   Semantic complexity: same diagnosis expressed with different phrasing",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "CLINICAL IMPACT", CL, Y0 + Inches(2.85), c=TEAL)
    _tm(s, CL + Inches(0.15), Y0 + Inches(3.13), CW, Inches(0.6), [
        "•   Manual coding costs US hospitals ~$25B/yr — automated prediction reduces errors",
        "•   Accurate codes drive reimbursement, epidemiology, and quality-of-care metrics",
    ], sz=11, c=BODY, sp=1.3)

    # ═══ 4  PRIOR WORK ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Prior Work", SW)
    _footer(s, SW, SH, 4)
    _tr(s, "fade")

    pw = [
        ("CAML / DR-CAML (Mullenbach et al., 2018)",
         "CNN + per-label attention over clinical text. Established that label-aware attention is critical for ICD coding."),
        ("PLM-ICD (Huang et al., 2022)",
         "Overlapping chunks + per-label attention over pretrained LM. Achieves SOTA on MIMIC-III with full-document coverage."),
        ("LAAT (Vu et al., 2020)",
         "BiLSTM encoder + label attention. Showed RNN-based models with label attention rival transformer approaches."),
        ("ClinicalBERT (Alsentzer et al., 2019)",
         "Pre-trained on MIMIC clinical notes; captures contextual semantics but truncates at 512 tokens."),
    ]
    for i, (h, b) in enumerate(pw):
        y = Y0 + Inches(0.05) + i * Inches(0.9)
        _r(s, CL, y, Inches(0.05), Inches(0.65), LBLUE)
        _t(s, CL + Inches(0.15), y, Inches(8.5), Inches(0.25), h, sz=12, c=BLUE, b=True)
        _t(s, CL + Inches(0.15), y + Inches(0.28), CW, Inches(0.45), b, sz=11, c=BODY)

    # ═══ 5  DATASET & BASE MODEL ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Dataset & Base Model (TF-IDF)", SW)
    _footer(s, SW, SH, 5)
    _tr(s, "fade")

    _sec(s, "MIMIC-IV CLINICAL DATABASE", CL, Y0)
    _tm(s, CL + Inches(0.15), Y0 + Inches(0.28), Inches(4.35), Inches(1.45), [
        "•  331,793 discharge notes from Beth Israel Deaconess",
        "•  122,304 admissions with notes + ICD-10 codes",
        "•  7,940 codes (freq ≥ 10); Top-50 for experiments",
        "•  Patient-level split (no leakage): Train 85,081 (70%), Val 18,371 (15%), Test 18,852 (15%)",
    ], sz=11, c=BODY, sp=1.3)

    # EDA plots — embedded outputs from Notebook 01 (re-extracted on each build)
    _img(s, ROOT / "data" / "note_length_dist.png",
         Inches(5.05), Y0 + Inches(0.05), width=Inches(4.65), height=Inches(1.05))
    _img(s, ROOT / "data" / "label_cardinality.png",
         Inches(5.05), Y0 + Inches(1.15), width=Inches(4.65), height=Inches(1.05))
    _img(s, ROOT / "data" / "code_freq_tail.png",
         Inches(5.05), Y0 + Inches(2.25), width=Inches(4.65), height=Inches(1.0))

    _sec(s, "BASE MODEL: TF-IDF + SGD (SHARED)", CL, Y0 + Inches(3.38), c=ACCENT)
    _tm(s, CL + Inches(0.15), Y0 + Inches(3.66), CW, Inches(1.45), [
        "•  50K-dim TF-IDF vectors (unigrams + bigrams, sublinear TF)",
        "•  OneVsRest SGDClassifier (log loss, L2 penalty, balanced class weights)",
        "•  Threshold tuning on validation → t = 0.53",
        "•  Micro-F1: 0.60 | Macro-F1: 0.57 | AUROC: 0.93",
        "•  Serves as the baseline both approaches aim to improve upon",
    ], sz=11, c=BODY, sp=1.25)

    # ═══ 6  CONTRIBUTIONS — TEAMMATE 1 ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Individual Contributions — Teammate 1", SW, sub="[Teammate 1 — edit name]")
    _footer(s, SW, SH, 6)
    _tr(s, "cover")

    _sec(s, "APPROACH 1: CHUNK-BERT + LABEL ATTENTION (PLM-ICD INSPIRED)", CL, Y0)
    _tm(s, CL + Inches(0.15), Y0 + Inches(0.28), CW, Inches(0.7), [
        "•  Designed chunking strategy: sliding window (6 × 512 tokens, stride 256) for full-document coverage",
        "•  Implemented per-label attention — 50 learnable query vectors attend across all chunk embeddings",
        "•  Bio_ClinicalBERT encoder with two-phase training (frozen → fine-tune last 2 layers)",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "TRAINING & OPTIMIZATION", CL, Y0 + Inches(1.15), c=ACCENT)
    _tm(s, CL + Inches(0.15), Y0 + Inches(1.43), CW, Inches(0.7), [
        "•  Sigmoid focal loss (gamma=2.0, alpha=0.25) to handle severe class imbalance",
        "•  Mixed-precision (fp16) + gradient accumulation (batch 4 × 8 = effective 32) on GCP L4 GPU",
        "•  ~8 hours training; validation F1 improved from 0.41 → 0.47 across training phases",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "ADDITIONAL CONTRIBUTIONS", CL, Y0 + Inches(2.3), c=TEAL)
    _tm(s, CL + Inches(0.15), Y0 + Inches(2.58), CW, Inches(0.9), [
        "•  MIMIC-IV data extraction pipeline & cohort construction (122K admissions)",
        "•  Text preprocessing: de-id stripping, normalization, TF-IDF feature engineering",
        "•  Built shared evaluation framework (threshold tuning, head/torso/tail analysis)",
    ], sz=11, c=BODY, sp=1.3)

    # ═══ 7  CONTRIBUTIONS — TEAMMATE 2 ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Individual Contributions — Teammate 2", SW, sub="[Teammate 2 — edit name]")
    _footer(s, SW, SH, 7)
    _tr(s, "cover")

    _sec(s, "APPROACH 2: BiLSTM-LAAT (LABEL ATTENTION OVER BiLSTM)", CL, Y0)
    _tm(s, CL + Inches(0.15), Y0 + Inches(0.28), CW, Inches(0.7), [
        "•  Implemented LAAT architecture: BiLSTM encoder + structured label attention mechanism",
        "•  Word-level tokenization with 50K vocabulary — processes up to 4,000 tokens natively",
        "•  Per-label attention with learned projection: Z = tanh(W·H^T), A = softmax(U·Z)",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "TRAINING & CALIBRATION", CL, Y0 + Inches(1.15), c=ACCENT)
    _tm(s, CL + Inches(0.15), Y0 + Inches(1.43), CW, Inches(0.7), [
        "•  10 epochs with focal loss, AdamW (lr=1e-3), batch size 32 on GCP L4 GPU",
        "•  Temperature calibration post-training — reduced ECE from 0.07 to 0.02",
        "•  Achieved best individual Micro-F1: 0.72 — significantly outperforms base model (+12 F1 pts)",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "ADDITIONAL CONTRIBUTIONS", CL, Y0 + Inches(2.3), c=TEAL)
    _tm(s, CL + Inches(0.15), Y0 + Inches(2.58), CW, Inches(0.9), [
        "•  Comparative evaluation across all models (per-label scatter analysis, bar charts)",
        "•  Production packaging: src/ Python module, FastAPI API, Docker containerization",
        "•  Pytest test suite covering data processing, model shapes, and API schemas",
    ], sz=11, c=BODY, sp=1.3)

    # ═══ 8  METHODS OVERVIEW ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Proposed Approaches — Overview", SW)
    _footer(s, SW, SH, 8)
    _tr(s, "push")

    _t(s, CL, Y0, CW, Inches(0.25),
       "Two approaches for Top-50 ICD-10 prediction, both addressing full-document processing:",
       sz=11, c=LGRAY)

    cards = [
        ("Approach 1", "Chunk-BERT + Label Attention", BLUE, [
            "PLM-ICD inspired architecture",
            "6 overlapping 512-token chunks → Bio_ClinicalBERT",
            "Concatenate all token embeddings across chunks",
            "Per-label attention: 50 learnable query vectors",
            "Two-phase training: frozen → fine-tune BERT",
        ]),
        ("Approach 2", "BiLSTM-LAAT", TEAL, [
            "LAAT (Vu et al., 2020) inspired architecture",
            "Word embeddings → Bidirectional LSTM encoder",
            "Structured label attention: W·H^T → tanh → softmax",
            "Processes up to 4,000 tokens natively (no chunking)",
            "Lightweight: 10 epochs, trains in ~30 min on L4",
        ]),
    ]
    for ci, (name, sub, hc, bul) in enumerate(cards):
        x = CL + ci * Inches(4.6)
        _r(s, x, Y0 + Inches(0.35), Inches(4.2), Inches(0.45), hc)
        _t(s, x + Inches(0.15), Y0 + Inches(0.37), Inches(3.9), Inches(0.42),
           name, sz=16, c=WHITE, b=True, va=MSO_ANCHOR.MIDDLE)
        _r(s, x, Y0 + Inches(0.8), Inches(4.2), Inches(2.85), LTBG)
        _t(s, x + Inches(0.15), Y0 + Inches(0.9), Inches(3.9), Inches(0.25),
           sub, sz=11, c=hc, b=True)
        _tm(s, x + Inches(0.15), Y0 + Inches(1.2), Inches(3.9), Inches(2.2),
            [f"•  {b}" for b in bul], sz=11, c=BODY, sp=1.35)

    _r(s, CL, Y0 + Inches(3.35), CW, Inches(0.38), LTBG)
    _t(s, CL + Inches(0.12), Y0 + Inches(3.38), CW - Inches(0.24), Inches(0.32),
       "Both approaches use per-label attention and process full documents — key advantages over truncation-based BERT.",
       sz=10, c=TEAL, b=True)

    # ═══ 9  APPROACH 1 DEEP DIVE ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Approach 1 — Chunk-BERT + Label Attention", SW)
    _footer(s, SW, SH, 9)
    _tr(s, "wipe")

    _sec(s, "CHUNKING + ENCODING", CL, Y0)
    _tm(s, CL + Inches(0.15), Y0 + Inches(0.28), Inches(4.8), Inches(0.85), [
        "1.  Tokenize full note → sliding window (512 tokens, stride 256)",
        "2.  Pad/truncate to max 6 chunks; mask padding during attention",
        "3.  Bio_ClinicalBERT encodes each chunk → concat all embeddings",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "LABEL ATTENTION", CL, Y0 + Inches(1.2), c=TEAL)
    _tm(s, CL + Inches(0.15), Y0 + Inches(1.48), Inches(4.8), Inches(0.7), [
        "•  50 learnable query vectors (R^768), Xavier-initialized",
        "•  Scores = H @ Q^T → softmax → weighted sum per label",
        "•  Per-label linear projection for final logits",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "TWO-PHASE TRAINING", CL, Y0 + Inches(2.3), c=ACCENT)
    _tm(s, CL + Inches(0.15), Y0 + Inches(2.58), Inches(4.8), Inches(1.0), [
        "Phase 1:  5 ep, lr=1e-3 — trains attention + classifier (BERT frozen)",
        "Phase 2:  3 ep, lr=2e-5 — unfreezes last 2 BERT layers",
        "Focal loss (gamma=2.0) · fp16 · batch 4×8 on GCP L4",
    ], sz=11, c=BODY, sp=1.3)

    _img(s, ROOT / "data" / "models" / "model_c" / "training_loss.png",
         Inches(5.3), Y0 + Inches(0.0), width=Inches(4.3), height=Inches(1.65))
    _img(s, ROOT / "data" / "models" / "model_c" / "training_f1.png",
         Inches(5.3), Y0 + Inches(1.7), width=Inches(4.3), height=Inches(1.65))

    # ═══ 10  APPROACH 2 DEEP DIVE ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Approach 2 — BiLSTM-LAAT", SW)
    _footer(s, SW, SH, 10)
    _tr(s, "wipe")

    _sec(s, "ARCHITECTURE (LAAT — Vu et al., 2020)", CL, Y0)
    _tm(s, CL + Inches(0.15), Y0 + Inches(0.28), Inches(4.8), Inches(0.85), [
        "1.  Word embedding layer (dim=200, vocab=50K) → BiLSTM (hidden=256)",
        "2.  Label attention: Z = tanh(W·H^T), A = softmax(U·Z)",
        "3.  Per-label FFN classifier → 50 sigmoid outputs",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "KEY DESIGN CHOICES", CL, Y0 + Inches(1.2), c=TEAL)
    _tm(s, CL + Inches(0.15), Y0 + Inches(1.48), Inches(4.8), Inches(0.7), [
        "•  Processes up to 4,000 word tokens — full document, no chunking",
        "•  Lightweight: ~11M parameters (vs 108M for BERT-based approach)",
        "•  Focal loss + temperature calibration (ECE: 0.07 → 0.02)",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "TRAINING", CL, Y0 + Inches(2.3), c=ACCENT)
    _tm(s, CL + Inches(0.15), Y0 + Inches(2.58), Inches(4.8), Inches(1.0), [
        "10 epochs, AdamW lr=1e-3, batch 32 on GCP L4 GPU",
        "Val F1 improved steadily: 0.67 → 0.72 (tuned threshold)",
        "Early stopping patience = 3; best at epoch 8",
    ], sz=11, c=BODY, sp=1.3)

    _img(s, ROOT / "data" / "models" / "model_d" / "training_loss.png",
         Inches(5.3), Y0 + Inches(0.0), width=Inches(4.3), height=Inches(1.65))
    _img(s, ROOT / "data" / "models" / "model_d" / "training_f1.png",
         Inches(5.3), Y0 + Inches(1.7), width=Inches(4.3), height=Inches(1.65))

    # ═══ 11  RESULTS TABLE ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Results — Test Set (Top-50 Codes)", SW)
    _footer(s, SW, SH, 11)
    _tr(s, "cover")

    _tbl(s, [["Metric", "Base Model\n(TF-IDF+SGD)", "Approach 1\n(Chunk+Attn)", "Approach 2\n(BiLSTM-LAAT)"],
             ["Micro-F1",        "0.60", "0.53", "0.72"],
             ["Macro-F1",        "0.57", "0.50", "0.66"],
             ["Micro-Precision", "0.49", "0.42", "0.72"],
             ["Micro-Recall",    "0.75", "0.72", "0.72"],
             ["Macro-AUPRC",     "0.57", "0.52", "0.68"],
             ["Micro-AUROC",     "0.93", "0.89", "0.95"],
             ["Threshold",       "0.53",  "0.63",  "0.30"]],
         [Inches(1.6), Inches(1.9), Inches(1.9), Inches(2.0)],
         Inches(0.6), Y0 + Inches(0.05), best_col=3)

    _r(s, CL, Y0 + Inches(3.2), CW, Inches(0.50), LTBG)
    _tm(s, CL + Inches(0.12), Y0 + Inches(3.23), CW - Inches(0.24), Inches(0.45), [
       "BiLSTM-LAAT achieves best results: +12 F1 pts over base model, with balanced precision (0.72) and recall (0.72).",
       "Approach 1 matches base model recall but needs more training epochs — BERT layers unfrozen only for 3 epochs.",
    ], sz=10, c=TEAL, b=True, sp=1.3)

    # ═══ 12  RESULTS — VISUAL ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Results — Model Comparison", SW)
    _footer(s, SW, SH, 12)
    _tr(s, "fade")

    img_path = ROOT / "data" / "models" / "three_model_comparison_bars.png"
    _img(s, img_path, Inches(0.3), Y0 + Inches(0.0), width=Inches(9.4), height=Inches(3.5))

    # ═══ 13  KEY TAKEAWAYS ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Key Takeaways", SW)
    _footer(s, SW, SH, 13)
    _tr(s, "push")

    tks = [
        ("BiLSTM-LAAT outperforms both BERT-based and TF-IDF approaches",
         "With only ~11M parameters, the BiLSTM achieves 0.72 F1 — processing 4K tokens natively without chunking overhead."),
        ("Full-document processing is critical for ICD coding",
         "Both proposed approaches process the complete note. Truncated ClinicalBERT (512 tokens) scored only 0.52 F1."),
        ("Label attention is the key architectural component",
         "Both approaches use per-label attention — each ICD code learns to attend to its most relevant tokens in the document."),
        ("Lightweight models can outperform large transformers",
         "BiLSTM-LAAT trains in ~30 min vs ~8 hrs for Chunk-BERT. Smaller models can be more practical for clinical deployment."),
    ]
    for i, (hd, bd) in enumerate(tks):
        y = Y0 + Inches(0.05) + i * Inches(0.85)
        _circ(s, CL, y + Inches(0.02), i + 1)
        _t(s, CL + Inches(0.48), y, Inches(8.0), Inches(0.25), hd, sz=13, c=BLUE, b=True)
        _t(s, CL + Inches(0.48), y + Inches(0.27), CW - Inches(0.5), Inches(0.45), bd, sz=11, c=BODY)

    # ═══ 14  DISCUSSION ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Discussion — Analysis & Insights", SW)
    _footer(s, SW, SH, 14)
    _tr(s, "wipe")

    _sec(s, "WHY DOES BiLSTM-LAAT OUTPERFORM CHUNK-BERT?", CL, Y0)
    _tm(s, CL + Inches(0.15), Y0 + Inches(0.28), CW, Inches(0.7), [
        "BiLSTM processes 4,000 tokens in a single pass — no chunking artifacts or cross-chunk boundary issues.",
        "Word-level tokenization preserves medical terminology directly; BERT's subword tokenizer fragments clinical terms.",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "WHY DOES TF-IDF REMAIN COMPETITIVE?", CL, Y0 + Inches(1.25), c=TEAL)
    _tm(s, CL + Inches(0.15), Y0 + Inches(1.53), CW, Inches(0.7), [
        "ICD codes are defined by precise medical terms — \"acute kidney failure\" maps directly to N17.9.",
        "TF-IDF bigrams capture these domain-specific phrases without needing contextual understanding.",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "CALIBRATION MATTERS FOR CLINICAL USE", CL, Y0 + Inches(2.5), c=ACCENT)
    _tm(s, CL + Inches(0.15), Y0 + Inches(2.78), CW, Inches(1.0), [
        "Temperature scaling on BiLSTM-LAAT reduced ECE from 0.07 to 0.02 — clinically reliable confidence scores.",
        "The gap between AUROC (0.95) and F1 (0.72) shows strong ranking ability — per-label threshold tuning could push F1 higher.",
    ], sz=11, c=BODY, sp=1.3)

    # ═══ 15  CHALLENGES ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Discussion — Challenges & Solutions", SW)
    _footer(s, SW, SH, 15)
    _tr(s, "push")

    ch = [
        ("Label-Space Scalability",
         "Full 7,940-label training was infeasible (100+ GPU-hrs). Focused on Top-50 codes covering 32.5% of labels."),
        ("Long-Document Processing",
         "BERT's 512-token limit loses 40-60% of each note. Solved with chunking (Approach 1) and BiLSTM (Approach 2)."),
        ("Severe Class Imbalance",
         "Rare codes appear in <0.1% of admissions. Addressed with focal loss (gamma=2.0) in both approaches."),
        ("Probability Calibration",
         "Raw model outputs were overconfident. Temperature scaling (T=0.63 for BiLSTM) reduced ECE to 0.02."),
    ]
    for i, (h, b) in enumerate(ch):
        y = Y0 + Inches(0.05) + i * Inches(0.9)
        _r(s, CL, y, Inches(0.05), Inches(0.65), ACCENT)
        _t(s, CL + Inches(0.15), y, Inches(4), Inches(0.25), h, sz=12, c=BLUE, b=True)
        _t(s, CL + Inches(0.15), y + Inches(0.28), CW, Inches(0.45), b, sz=11, c=BODY)

    # ═══ 16  FUTURE WORK ═══
    s = prs.slides.add_slide(BL)
    _hdr(s, "Future Work & Next Steps", SW)
    _footer(s, SW, SH, 16)
    _tr(s, "fade")

    _sec(s, "ENSEMBLE — IMMEDIATE NEXT STEP", CL, Y0)
    _tm(s, CL + Inches(0.15), Y0 + Inches(0.28), CW, Inches(0.85), [
        "•  Weighted probability blend of TF-IDF + BiLSTM-LAAT with tuned threshold",
        "•  The models make different types of errors — lexical vs. sequential-contextual",
        "•  Preliminary analysis suggests 3-5% F1 improvement is achievable over best individual model",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "INTERACTIVE DEMO (STREAMLIT + FastAPI)", CL, Y0 + Inches(1.35), c=TEAL)
    _tm(s, CL + Inches(0.15), Y0 + Inches(1.63), CW, Inches(0.85), [
        "•  Streamlit web UI: paste a discharge summary → predicted ICD-10 codes with confidence scores",
        "•  FastAPI backend with attention-based evidence highlighting per predicted code",
        "•  Docker-containerized for reproducible deployment",
    ], sz=11, c=BODY, sp=1.3)

    _sec(s, "FURTHER EXTENSIONS", CL, Y0 + Inches(2.7), c=ACCENT)
    _tm(s, CL + Inches(0.15), Y0 + Inches(2.98), CW, Inches(1.0), [
        "•  Scale to Top-500 codes — more realistic clinical deployment scenario",
        "•  Per-label threshold tuning — global threshold is suboptimal for rare codes",
        "•  Clinical-Longformer (4,096 tokens) to combine transformer power with full-document coverage",
    ], sz=11, c=BODY, sp=1.3)

    # ═══ 17  THANK YOU ═══
    s = prs.slides.add_slide(BL)
    _r(s, 0, 0, SW, SH, BLUE)
    _r(s, 0, Inches(2.6), SW, Inches(0.04), ACCENT)
    _t(s, Inches(0.8), Inches(1.0), Inches(8.4), Inches(0.9),
       "Thank You!", sz=38, c=WHITE, b=True, a=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    _t(s, Inches(0.8), Inches(2.85), Inches(8.4), Inches(0.4),
       "Questions?", sz=20, c=RGBColor(0x99, 0xAA, 0xBB), a=PP_ALIGN.CENTER)
    _t(s, Inches(0.8), Inches(3.6), Inches(8.4), Inches(0.3),
       "ICD-10 Code Prediction  ·  NLP Final Project  ·  Spring 2025",
       sz=12, c=RGBColor(0x88, 0x99, 0xAA), a=PP_ALIGN.CENTER)
    _t(s, Inches(0.8), Inches(4.05), Inches(8.4), Inches(0.35),
       "Devasai Sunder Tangella    ·    [Partner Name]",
       sz=15, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    _tr(s, "fade")

    prs.save(str(out))


def main():
    root = Path(__file__).resolve().parents[1]
    build(root / "template_clean.pptx", root / "Final_Project_Presentation.pptx")
    print("Done → Final_Project_Presentation.pptx  (17 slides)")


if __name__ == "__main__":
    main()
